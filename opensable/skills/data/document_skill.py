"""
Document creation skill for Open-Sable
Creates Word, Excel, PDF, and PowerPoint files using pure Python libraries.
Cross-platform: works on Windows, macOS, and Linux without LibreOffice.
"""

import logging
import json
from typing import Dict, Any, List, Optional
from pathlib import Path
from datetime import datetime

logger = logging.getLogger(__name__)

# Default output directory
_DEFAULT_DIR = Path.home() / "Documents" / "SableDocs"


def _ensure_dir(path: Path) -> Path:
    """Ensure the parent directory exists; return the path."""
    path.parent.mkdir(parents=True, exist_ok=True)
    return path


def _resolve_output(filename: str, ext: str, output_dir: Optional[str] = None) -> Path:
    """Build a full output path, adding extension if missing."""
    base = Path(output_dir) if output_dir else _DEFAULT_DIR
    base.mkdir(parents=True, exist_ok=True)
    if not filename.lower().endswith(ext):
        filename += ext
    return base / filename


class DocumentSkill:
    """Create Word (.docx), Excel (.xlsx), PDF, and PowerPoint (.pptx) documents."""

    def __init__(self, config):
        self.config = config
        self._available = {"docx": False, "xlsx": False, "pdf": False, "pptx": False}

    async def initialize(self) -> bool:
        """Check which document libraries are installed."""
        try:
            import docx  # noqa: F401
            self._available["docx"] = True
        except ImportError:
            logger.warning("python-docx not installed,  Word creation disabled")

        try:
            import openpyxl  # noqa: F401
            self._available["xlsx"] = True
        except ImportError:
            logger.warning("openpyxl not installed,  Excel creation disabled")

        try:
            from reportlab.lib.pagesizes import A4  # noqa: F401
            self._available["pdf"] = True
        except ImportError:
            logger.warning("reportlab not installed,  PDF creation disabled")

        try:
            from pptx import Presentation  # noqa: F401
            self._available["pptx"] = True
        except ImportError:
            logger.warning("python-pptx not installed,  PowerPoint creation disabled")

        avail = [k for k, v in self._available.items() if v]
        logger.info(f"DocumentSkill initialized,  available formats: {avail or 'none'}")
        return bool(avail)

    # ------------------------------------------------------------------
    # Word (.docx)
    # ------------------------------------------------------------------

    async def create_word(
        self,
        filename: str = "document.docx",
        title: str = "",
        content: str = "",
        paragraphs: Optional[List[str]] = None,
        table_data: Optional[List[List[str]]] = None,
        output_dir: Optional[str] = None,
    ) -> Dict[str, Any]:
        """Create a Word document.

        Args:
            filename: Output file name.
            title: Document title (added as Heading 1).
            content: Body text (if paragraphs not provided).
            paragraphs: List of paragraph strings.
            table_data: 2-D list of strings for a table (first row = headers).
            output_dir: Override output directory.
        """
        if not self._available["docx"]:
            return {"success": False, "error": "python-docx is not installed. Run: pip install python-docx"}

        try:
            from docx import Document
            from docx.shared import Pt, Inches

            doc = Document()

            if title:
                doc.add_heading(title, level=1)

            # Body text
            body = paragraphs or ([content] if content else [])
            for para in body:
                doc.add_paragraph(para)

            # Optional table
            if table_data and len(table_data) >= 1:
                cols = len(table_data[0])
                table = doc.add_table(rows=1, cols=cols, style="Light Grid Accent 1")
                # Header row
                for i, header in enumerate(table_data[0]):
                    table.rows[0].cells[i].text = str(header)
                # Data rows
                for row_data in table_data[1:]:
                    row = table.add_row()
                    for i, cell in enumerate(row_data):
                        if i < cols:
                            row.cells[i].text = str(cell)

            out = _resolve_output(filename, ".docx", output_dir)
            doc.save(str(out))
            logger.info(f"Created Word document: {out}")
            return {"success": True, "path": str(out), "format": "docx"}

        except Exception as e:
            logger.error(f"Word creation failed: {e}")
            return {"success": False, "error": str(e)}

    # ------------------------------------------------------------------
    # Excel (.xlsx)
    # ------------------------------------------------------------------

    async def create_spreadsheet(
        self,
        filename: str = "spreadsheet.xlsx",
        sheets: Optional[Dict[str, List[List[Any]]]] = None,
        data: Optional[List[List[Any]]] = None,
        headers: Optional[List[str]] = None,
        output_dir: Optional[str] = None,
    ) -> Dict[str, Any]:
        """Create an Excel spreadsheet.

        Args:
            filename: Output file name.
            sheets: Dict mapping sheet names to 2-D data arrays.
            data: Simple 2-D array for a single "Sheet1".
            headers: Column headers (prepended to data).
            output_dir: Override output directory.
        """
        if not self._available["xlsx"]:
            return {"success": False, "error": "openpyxl is not installed. Run: pip install openpyxl"}

        try:
            from openpyxl import Workbook
            from openpyxl.styles import Font, PatternFill, Alignment
            from openpyxl.utils import get_column_letter

            wb = Workbook()

            # Build sheet map
            if sheets:
                sheet_map = sheets
            elif data:
                rows = data
                if headers:
                    rows = [headers] + rows
                sheet_map = {"Sheet1": rows}
            else:
                sheet_map = {"Sheet1": [["(empty spreadsheet)"]]}

            first = True
            for sheet_name, rows in sheet_map.items():
                if first:
                    ws = wb.active
                    ws.title = sheet_name
                    first = False
                else:
                    ws = wb.create_sheet(title=sheet_name)

                for r_idx, row in enumerate(rows, 1):
                    for c_idx, value in enumerate(row, 1):
                        cell = ws.cell(row=r_idx, column=c_idx, value=value)
                        # Bold header row
                        if r_idx == 1:
                            cell.font = Font(bold=True)
                            cell.fill = PatternFill(start_color="D9E1F2", end_color="D9E1F2", fill_type="solid")
                            cell.alignment = Alignment(horizontal="center")

                # Auto-width columns
                for col_idx in range(1, ws.max_column + 1):
                    max_len = 0
                    col_letter = get_column_letter(col_idx)
                    for row in ws.iter_rows(min_col=col_idx, max_col=col_idx):
                        for cell in row:
                            if cell.value:
                                max_len = max(max_len, len(str(cell.value)))
                    ws.column_dimensions[col_letter].width = min(max_len + 4, 50)

            out = _resolve_output(filename, ".xlsx", output_dir)
            wb.save(str(out))
            logger.info(f"Created spreadsheet: {out}")
            return {"success": True, "path": str(out), "format": "xlsx", "sheets": list(sheet_map.keys())}

        except Exception as e:
            logger.error(f"Spreadsheet creation failed: {e}")
            return {"success": False, "error": str(e)}

    # ------------------------------------------------------------------
    # PDF
    # ------------------------------------------------------------------

    async def create_pdf(
        self,
        filename: str = "document.pdf",
        title: str = "",
        content: str = "",
        paragraphs: Optional[List[str]] = None,
        table_data: Optional[List[List[str]]] = None,
        output_dir: Optional[str] = None,
    ) -> Dict[str, Any]:
        """Create a PDF document.

        Args:
            filename: Output file name.
            title: Document title (large centered text at top).
            content: Body text (if paragraphs not provided).
            paragraphs: List of paragraph strings.
            table_data: 2-D list for a table (first row = headers).
            output_dir: Override output directory.
        """
        if not self._available["pdf"]:
            return {"success": False, "error": "reportlab is not installed. Run: pip install reportlab"}

        try:
            from reportlab.lib.pagesizes import A4
            from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
            from reportlab.lib.units import cm
            from reportlab.lib.colors import HexColor
            from reportlab.platypus import (
                SimpleDocTemplate, Paragraph, Spacer, Table, TableStyle, PageBreak,
            )

            out = _resolve_output(filename, ".pdf", output_dir)
            doc = SimpleDocTemplate(
                str(out), pagesize=A4,
                leftMargin=2 * cm, rightMargin=2 * cm,
                topMargin=2 * cm, bottomMargin=2 * cm,
            )

            styles = getSampleStyleSheet()
            story: list = []

            if title:
                title_style = ParagraphStyle(
                    "DocTitle", parent=styles["Title"],
                    fontSize=22, spaceAfter=20,
                )
                story.append(Paragraph(title, title_style))

            body = paragraphs or ([content] if content else [])
            for para in body:
                # Replace newlines with <br/> for Platypus
                safe = para.replace("&", "&amp;").replace("<", "&lt;").replace(">", "&gt;")
                safe = safe.replace("\n", "<br/>")
                story.append(Paragraph(safe, styles["BodyText"]))
                story.append(Spacer(1, 8))

            if table_data and len(table_data) >= 1:
                t = Table(table_data, repeatRows=1)
                t.setStyle(TableStyle([
                    ("BACKGROUND", (0, 0), (-1, 0), HexColor("#4472C4")),
                    ("TEXTCOLOR", (0, 0), (-1, 0), HexColor("#FFFFFF")),
                    ("FONTNAME", (0, 0), (-1, 0), "Helvetica-Bold"),
                    ("ALIGN", (0, 0), (-1, -1), "LEFT"),
                    ("GRID", (0, 0), (-1, -1), 0.5, HexColor("#CCCCCC")),
                    ("ROWBACKGROUNDS", (0, 1), (-1, -1), [HexColor("#F2F2F2"), HexColor("#FFFFFF")]),
                    ("TOPPADDING", (0, 0), (-1, -1), 4),
                    ("BOTTOMPADDING", (0, 0), (-1, -1), 4),
                ]))
                story.append(Spacer(1, 12))
                story.append(t)

            doc.build(story)
            logger.info(f"Created PDF: {out}")
            return {"success": True, "path": str(out), "format": "pdf"}

        except Exception as e:
            logger.error(f"PDF creation failed: {e}")
            return {"success": False, "error": str(e)}

    # ------------------------------------------------------------------
    # PowerPoint (.pptx)
    # ------------------------------------------------------------------

    async def create_presentation(
        self,
        filename: str = "presentation.pptx",
        title: str = "",
        subtitle: str = "",
        slides: Optional[List[Dict[str, Any]]] = None,
        output_dir: Optional[str] = None,
    ) -> Dict[str, Any]:
        """Create a PowerPoint presentation.

        Args:
            filename: Output file name.
            title: Title slide heading.
            subtitle: Title slide subtitle.
            slides: List of dicts, each with:
                - title (str): Slide heading
                - content (str): Body text
                - bullets (list[str]): Bullet points
                - image (str): Path to image file to embed
                - layout (str): "title", "content", "bullets", "blank"
            output_dir: Override output directory.
        """
        if not self._available["pptx"]:
            return {"success": False, "error": "python-pptx is not installed. Run: pip install python-pptx"}

        try:
            from pptx import Presentation
            from pptx.util import Inches, Pt
            from pptx.enum.text import PP_ALIGN

            prs = Presentation()

            # Title slide
            if title:
                layout = prs.slide_layouts[0]  # Title slide
                slide = prs.slides.add_slide(layout)
                slide.shapes.title.text = title
                if subtitle:
                    slide.placeholders[1].text = subtitle

            # Content slides
            for slide_def in (slides or []):
                s_title = slide_def.get("title", "")
                s_content = slide_def.get("content", "")
                s_bullets = slide_def.get("bullets", [])
                s_image = slide_def.get("image", "")
                s_layout = slide_def.get("layout", "content")

                if s_layout == "blank":
                    slide = prs.slides.add_slide(prs.slide_layouts[6])
                elif s_layout == "title":
                    slide = prs.slides.add_slide(prs.slide_layouts[0])
                    slide.shapes.title.text = s_title
                    if s_content:
                        slide.placeholders[1].text = s_content
                    continue
                else:
                    # Title + Content layout
                    slide = prs.slides.add_slide(prs.slide_layouts[1])
                    slide.shapes.title.text = s_title

                    body = slide.placeholders[1]
                    tf = body.text_frame
                    tf.clear()

                    if s_bullets:
                        for i, bullet in enumerate(s_bullets):
                            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                            p.text = bullet
                            p.level = 0
                    elif s_content:
                        tf.paragraphs[0].text = s_content

                # Add image if provided
                if s_image and Path(s_image).exists():
                    slide.shapes.add_picture(s_image, Inches(1), Inches(2), width=Inches(4))

            out = _resolve_output(filename, ".pptx", output_dir)
            prs.save(str(out))
            logger.info(f"Created presentation: {out}")
            return {"success": True, "path": str(out), "format": "pptx", "slide_count": len(prs.slides)}

        except Exception as e:
            logger.error(f"Presentation creation failed: {e}")
            return {"success": False, "error": str(e)}

    # ------------------------------------------------------------------
    # Read / convert existing documents
    # ------------------------------------------------------------------

    async def read_document(self, file_path: str) -> Dict[str, Any]:
        """Read text content from a Word, Excel, or PDF file.

        Useful for extracting content before processing or summarizing.
        """
        path = Path(file_path)
        if not path.exists():
            return {"success": False, "error": f"File not found: {file_path}"}

        ext = path.suffix.lower()

        try:
            if ext == ".docx":
                return await self._read_docx(path)
            elif ext in (".xlsx", ".xls"):
                return await self._read_xlsx(path)
            elif ext == ".pdf":
                return await self._read_pdf(path)
            elif ext == ".pptx":
                return await self._read_pptx(path)
            elif ext in (".txt", ".md", ".csv", ".json", ".xml", ".html"):
                text = path.read_text(encoding="utf-8", errors="replace")
                return {"success": True, "text": text, "format": ext.lstrip(".")}
            else:
                return {"success": False, "error": f"Unsupported format: {ext}"}
        except Exception as e:
            return {"success": False, "error": str(e)}

    async def _read_docx(self, path: Path) -> Dict[str, Any]:
        from docx import Document
        doc = Document(str(path))
        paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
        return {"success": True, "text": "\n\n".join(paragraphs), "format": "docx", "paragraphs": len(paragraphs)}

    async def _read_xlsx(self, path: Path) -> Dict[str, Any]:
        from openpyxl import load_workbook
        wb = load_workbook(str(path), read_only=True, data_only=True)
        sheets = {}
        for name in wb.sheetnames:
            ws = wb[name]
            rows = []
            for row in ws.iter_rows(values_only=True):
                rows.append([str(c) if c is not None else "" for c in row])
            sheets[name] = rows
        wb.close()
        return {"success": True, "sheets": sheets, "format": "xlsx", "sheet_names": list(sheets.keys())}

    async def _read_pdf(self, path: Path) -> Dict[str, Any]:
        """Read PDF text. Tries PyMuPDF first, falls back to pdfplumber."""
        text = ""
        try:
            import fitz  # PyMuPDF
            doc = fitz.open(str(path))
            pages = []
            for page in doc:
                pages.append(page.get_text())
            doc.close()
            text = "\n\n".join(pages)
        except ImportError:
            try:
                import pdfplumber
                with pdfplumber.open(str(path)) as pdf:
                    pages = [p.extract_text() or "" for p in pdf.pages]
                text = "\n\n".join(pages)
            except ImportError:
                return {"success": False, "error": "Install PyMuPDF or pdfplumber to read PDFs: pip install PyMuPDF"}

        return {"success": True, "text": text, "format": "pdf", "pages": len(pages)}

    async def _read_pptx(self, path: Path) -> Dict[str, Any]:
        from pptx import Presentation
        prs = Presentation(str(path))
        slides_text = []
        for i, slide in enumerate(prs.slides, 1):
            parts = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    parts.append(shape.text_frame.text)
            slides_text.append(f"--- Slide {i} ---\n" + "\n".join(parts))
        return {"success": True, "text": "\n\n".join(slides_text), "format": "pptx", "slide_count": len(slides_text)}

    # ------------------------------------------------------------------
    # Open document in default application
    # ------------------------------------------------------------------

    async def open_document(self, file_path: str) -> Dict[str, Any]:
        """Open a document with the system's default application.
        Cross-platform: uses xdg-open (Linux), open (macOS), start (Windows).
        """
        import platform
        import subprocess

        path = Path(file_path)
        if not path.exists():
            return {"success": False, "error": f"File not found: {file_path}"}

        system = platform.system()
        try:
            if system == "Darwin":
                subprocess.Popen(["open", str(path)])
            elif system == "Windows":
                subprocess.Popen(["cmd", "/c", "start", "", str(path)], shell=False)
            else:
                subprocess.Popen(["xdg-open", str(path)])
            return {"success": True, "opened": str(path), "system": system}
        except Exception as e:
            return {"success": False, "error": str(e)}
